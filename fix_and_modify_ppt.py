import zipfile
import shutil
import os
from pptx import Presentation
from pptx.util import Pt

# Paths
src = 'C:/Users/zero/Desktop/校园导航系统_作品演示PPT.pptx'
temp_dir = 'C:/Users/zero/.openclaw/workspace/temp_pptx'
fixed_path = 'C:/Users/zero/.openclaw/workspace/校园导航系统_作品演示PPT_fixed.pptx'
output_path = 'C:/Users/zero/.openclaw/workspace/校园导航系统_作品演示PPT_修改版.pptx'

# Step 1: Extract and fix slide10.xml
if os.path.exists(temp_dir):
    shutil.rmtree(temp_dir)
os.makedirs(temp_dir)

with zipfile.ZipFile(src, 'r') as zin:
    zin.extractall(temp_dir)

# Read and fix slide10.xml
slide10_path = os.path.join(temp_dir, 'ppt/slides/slide10.xml')
with open(slide10_path, 'rb') as f:
    content = f.read()

# The corruption: '1' was replaced with 'O(n\xc2\xb2)' and '0' with 'O(1)'
# We need to reverse this. But be careful about order.
# Since O(1) contains '1', if we replace 'O(n\xc2\xb2)' first, we might affect O(1) -> no, O(1) doesn't contain O(n\xc2\xb2)
# But O(n\xc2\xb2) doesn't contain O(1) either.
# However, the original replacements might have been done in a specific order.
# Let's just do global replacements.

content_str = content.decode('utf-8')

# Replace the corrupted patterns back to numbers
# Pattern: O(n²) was used to replace "1"
# Pattern: O(1) was used to replace "0"

# But we need to be careful - in the actual slide content, O(n²) and O(1) are valid notations
# Since we're going to rewrite the content anyway, let's just fix the XML markup

# First, let's see what the actual text content says
import re
# Find all text runs
text_runs = re.findall(r'<a:t>([^<]*)</a:t>', content_str)
print("Text runs in slide10:")
for t in text_runs:
    print(f"  '{t}'")

# Do the replacement
content_str = content_str.replace('O(n²)', '1')
content_str = content_str.replace('O(1)', '0')

with open(slide10_path, 'w', encoding='utf-8') as f:
    f.write(content_str)

# Step 2: Repackage the PPTX
if os.path.exists(fixed_path):
    os.remove(fixed_path)

with zipfile.ZipFile(fixed_path, 'w', zipfile.ZIP_DEFLATED) as zout:
    for root, dirs, files in os.walk(temp_dir):
        for file in files:
            file_path = os.path.join(root, file)
            arcname = os.path.relpath(file_path, temp_dir)
            zout.write(file_path, arcname)

print(f"Fixed PPTX saved to: {fixed_path}")

# Step 3: Load and inspect
prs = Presentation(fixed_path)
print(f"\nTotal slides: {len(prs.slides)}")

for i, slide in enumerate(prs.slides):
    print(f"\n--- Slide {i+1} ---")
    for shape in slide.shapes:
        if shape.has_text_frame:
            text = shape.text_frame.text.strip()
            if text:
                print(f"  {shape.name}: {text[:100]}")

# Clean up temp dir
shutil.rmtree(temp_dir)
