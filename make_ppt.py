from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from copy import deepcopy
import os

# Load template
prs = Presentation('C:/Users/zero/.openclaw/workspace/template.pptx')

# Helper function to replace text in shapes
def replace_text_in_shape(shape, old_text, new_text):
    if shape.has_text_frame:
        for paragraph in shape.text_frame.paragraphs:
            for run in paragraph.runs:
                if old_text in run.text:
                    run.text = run.text.replace(old_text, new_text)
        # Also check the full text
        full_text = shape.text_frame.text
        if old_text in full_text:
            # Clear and rewrite
            shape.text_frame.clear()
            p = shape.text_frame.paragraphs[0]
            p.text = full_text.replace(old_text, new_text)

# Helper to set text preserving formatting
def set_shape_text(shape, new_text):
    if shape.has_text_frame:
        # Try to preserve formatting by replacing text in runs
        old_text = shape.text_frame.text
        if old_text:
            # Find first run with text
            for paragraph in shape.text_frame.paragraphs:
                for run in paragraph.runs:
                    if run.text:
                        run.text = new_text
                        return
        # Fallback: clear and set
        shape.text_frame.clear()
        p = shape.text_frame.paragraphs[0]
        p.text = new_text

print(f"Total slides: {len(prs.slides)}")

# Slide 1: Cover
slide1 = prs.slides[0]
for shape in slide1.shapes:
    if shape.has_text_frame:
        text = shape.text_frame.text.strip()
        print(f"Slide 1 shape: {shape.name} -> {text[:50] if text else '(empty)'}")
        if "学生成绩管理系统" in text:
            set_shape_text(shape, "增强版校园导航系统")
        elif "2501123121" in text:
            set_shape_text(shape, "汇报人学号：202501150111\n汇报人姓名：李泽锐")
        elif "2026-5-27" in text:
            set_shape_text(shape, "时间：2026-06-09")

# Slide 2: Table of Contents
slide2 = prs.slides[1]
for shape in slide2.shapes:
    if shape.has_text_frame:
        text = shape.text_frame.text.strip()
        if "子模的详细设计与实现解说" in text:
            set_shape_text(shape, "子模块的详细设计与实现解说")

# Slide 3: Section 1 divider
slide3 = prs.slides[2]
for shape in slide3.shapes:
    if shape.has_text_frame:
        text = shape.text_frame.text.strip()
        print(f"Slide 3 shape: {shape.name} -> {text[:50] if text else '(empty)'}")

# Slide 4: Content - System intro
slide4 = prs.slides[3]
for shape in slide4.shapes:
    if shape.has_text_frame:
        text = shape.text_frame.text.strip()
        print(f"Slide 4 shape: {shape.name} -> {text[:50] if text else '(empty)'}")
        if text == "系统简介":
            # Find the next shape with "。。。" and replace
            pass
        elif "。。。" in text and len(text) < 10:
            set_shape_text(shape, "基于Dijkstra算法和邻接矩阵的增强版校园导航系统，支持校园节点管理、路径规划、最短路径查询、数据持久化等功能。")
        elif "展示关键结构体定义" in text:
            set_shape_text(shape, " CampusMap结构体封装节点名称、邻接矩阵、节点数、路径数")
        elif "代码量与复杂度" in text:
            set_shape_text(shape, "约450行C语言代码，12个功能模块")
        elif "数据结构选型" in text:
            set_shape_text(shape, "邻接矩阵（二维数组）存储图结构，支持O(1)随机访问")
        elif "功能模块个数" in text:
            set_shape_text(shape, "12个功能模块：节点/路径增删查、最短路径、文件读写、图形化展示")

# Slide 5-6: Architecture diagram
slide5 = prs.slides[4]
for shape in slide5.shapes:
    if shape.has_text_frame:
        text = shape.text_frame.text.strip()
        print(f"Slide 5 shape: {shape.name} -> {text[:80] if text else '(empty)'}")
        if "系统架构图" in text and len(text) > 20:
            set_shape_text(shape, "系统架构图\n主菜单控制模块 → 节点管理模块 → 路径管理模块 → 查询模块 → 地图维护模块 → 数据持久化模块\n\n核心交互：主菜单通过switch-case调度12个功能函数，各模块共享全局CampusMap结构体变量。")
        elif "系统功能架构图" in text:
            set_shape_text(shape, "系统功能架构图\n主菜单 → 添加节点/路径 → 查询节点/矩阵/路径 → 删除节点/路径 → 保存/加载地图 → 图形化展示")

slide6 = prs.slides[5]
for shape in slide6.shapes:
    if shape.has_text_frame:
        text = shape.text_frame.text.strip()
        if "系统架构图" in text and len(text) > 20:
            set_shape_text(shape, "核心数据结构：\nstruct CampusMap {\n  char nodes[MAX][50];  // 节点名称\n  int graph[MAX][MAX];  // 邻接矩阵\n  int nodeNum;          // 节点总数\n  int edgeNum;          // 路径总数\n};")

# Slide 7: Section 2 divider
slide7 = prs.slides[6]
for shape in slide7.shapes:
    if shape.has_text_frame:
        text = shape.text_frame.text.strip()
        print(f"Slide 7 shape: {shape.name} -> {text[:50] if text else '(empty)'}")

# Slide 8: Demo
slide8 = prs.slides[7]
for shape in slide8.shapes:
    if shape.has_text_frame:
        text = shape.text_frame.text.strip()
        print(f"Slide 8 shape: {shape.name} -> {text[:80] if text else '(empty)'}")
        if "系统功能演示与测试用例" in text and len(text) > 20:
            set_shape_text(shape, "系统功能演示与测试用例\n\n1. 启动系统：输入口令1进入主菜单\n2. 添加节点：图书馆、教学楼、食堂、宿舍\n3. 添加路径：设置各节点间路径长度\n4. 查询最短路径：图书馆→宿舍，自动计算最优路线\n5. 查看邻接矩阵：直观展示图结构\n6. 保存/加载地图：数据持久化到campus_map.txt")

# Slide 9: Section 3 divider
slide9 = prs.slides[8]
for shape in slide9.shapes:
    if shape.has_text_frame:
        text = shape.text_frame.text.strip()
        print(f"Slide 9 shape: {shape.name} -> {text[:50] if text else '(empty)'}")

# Slide 10: Module overview
slide10 = prs.slides[9]
for shape in slide10.shapes:
    if shape.has_text_frame:
        text = shape.text_frame.text.strip()
        print(f"Slide 10 shape: {shape.name} -> {text[:50] if text else '(empty)'}")
        if "模块一" in text:
            set_shape_text(shape, "模块一\n节点与路径管理")
        elif "模块二" in text:
            set_shape_text(shape, "模块二\n地图查询与展示")
        elif "模块三" in text:
            set_shape_text(shape, "模块三\n最短路径计算")
        elif "模块四" in text:
            set_shape_text(shape, "模块四\n数据持久化")
        elif "插入新记录" in text:
            set_shape_text(shape, "添加节点/路径")
        elif "删除记录" in text:
            set_shape_text(shape, "删除节点/路径")
        elif "修改记录" in text:
            set_shape_text(shape, "查询节点/矩阵")
        elif "显示记录" in text:
            set_shape_text(shape, "保存/加载地图")

# Slide 11: Algorithm analysis
slide11 = prs.slides[10]
for shape in slide11.shapes:
    if shape.has_text_frame:
        text = shape.text_frame.text.strip()
        print(f"Slide 11 shape: {shape.name} -> {text[:80] if text else '(empty)'}")
        if "模块一" in text:
            set_shape_text(shape, "Dijkstra算法详解")
        elif "算法逻辑与复杂度分析" in text:
            set_shape_text(shape, "算法逻辑与复杂度分析（Dijkstra最短路径算法）\n\n算法步骤：\n1. 初始化dist数组、visit数组、pre数组\n2. 选取未访问节点中距离最小的节点u\n3. 对u的所有邻接节点执行松弛操作\n4. 重复直至所有节点访问完毕\n\n时间复杂度：O(n²) — 邻接矩阵实现\n空间复杂度：O(n) — dist、pre、visit数组\n\n适用条件：非负权图，完全满足本系统路径长度为正的设定")

# Slide 12: Code details
slide12 = prs.slides[11]
for shape in slide12.shapes:
    if shape.has_text_frame:
        text = shape.text_frame.text.strip()
        print(f"Slide 12 shape: {shape.name} -> {text[:80] if text else '(empty)'}")
        if "作品简要介绍" in text:
            set_shape_text(shape, "关键代码细节解析")
        elif "关键代码细节解析" in text:
            set_shape_text(shape, "关键代码细节解析（指针/内存/函数调用）\n\n1. 结构体封装：\n   struct CampusMap {\n     char nodes[MAX][50];\n     int graph[MAX][MAX];\n     int nodeNum, edgeNum;\n   };\n\n2. 邻接矩阵初始化：\n   graph[i][j] = (i==j) ? 0 : INF;\n\n3. 路径回溯（递归）：\n   void showPath(int end, int pre[]) {\n     if (pre[end] != end) {\n       showPath(pre[end], pre);\n       printf(\" -> \");\n     }\n     printf(\"%s\", nodes[end]);\n   }\n\n4. 文件持久化：\n   fprintf(fp, \"%d\\n\", nodeNum);\n   // 写入节点名称 + 邻接矩阵")

# Save the new presentation
output_path = 'C:/Users/zero/.openclaw/workspace/增强版校园导航系统_作品展示.pptx'
prs.save(output_path)
print(f"\nSaved to: {output_path}")
