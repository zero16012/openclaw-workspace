from pptx import Presentation
from pptx.util import Pt

# Load the working template
prs = Presentation('C:/Users/zero/.openclaw/workspace/template.pptx')

def clear_and_set_text(shape, new_text):
    """Clear all text in shape and set new text"""
    if not shape.has_text_frame:
        return
    
    tf = shape.text_frame
    
    # Save formatting from first paragraph's first run if exists
    saved_font_size = None
    saved_font_name = None
    saved_bold = None
    
    try:
        if tf.paragraphs:
            first_p = tf.paragraphs[0]
            if first_p.runs:
                first_run = first_p.runs[0]
                saved_font_size = first_run.font.size
                saved_font_name = first_run.font.name
                saved_bold = first_run.font.bold
    except:
        pass
    
    # Clear all text
    tf.clear()
    
    # Add new text
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = new_text
    
    # Restore formatting
    if saved_font_size:
        run.font.size = saved_font_size
    if saved_font_name:
        run.font.name = saved_font_name
    if saved_bold is not None:
        run.font.bold = saved_bold

# ========== SLIDE 1: Cover ==========
slide = prs.slides[0]
for shape in slide.shapes:
    if not shape.has_text_frame:
        continue
    text = shape.text_frame.text.strip()
    if not text:
        continue
    
    if "学生成绩管理系统" in text:
        clear_and_set_text(shape, "增强版校园导航系统")
    elif "2501123121" in text:
        clear_and_set_text(shape, "汇报人学号：202501150111\n汇报人姓名：李泽锐")
    elif "2026-5-27" in text:
        clear_and_set_text(shape, "时间：2026-06-09")

# ========== SLIDE 2: TOC ==========
slide = prs.slides[1]
for shape in slide.shapes:
    if not shape.has_text_frame:
        continue
    text = shape.text_frame.text.strip()
    if "子模的详细设计与实现解说" in text:
        clear_and_set_text(shape, "子模块的详细设计与实现解说")

# ========== SLIDE 4: Content - System intro ==========
slide = prs.slides[3]
for shape in slide.shapes:
    if not shape.has_text_frame:
        continue
    text = shape.text_frame.text.strip()
    
    if shape.name.startswith("Docer"):
        if "系统简介" in text and len(text) < 10:
            clear_and_set_text(shape, "系统简介")
        elif "。。。" in text:
            clear_and_set_text(shape, "基于C语言+Dijkstra算法+邻接矩阵的增强版校园导航系统")
        elif "代码量与复杂度" in text and len(text) < 15:
            clear_and_set_text(shape, "代码量与复杂度")
        elif "展示关键结构体定义" in text:
            clear_and_set_text(shape, "约450行C代码，16个功能函数")
        elif "数据结构选型" in text and len(text) < 10:
            clear_and_set_text(shape, "数据结构选型")
        elif " CampusMap结构体" in text:
            clear_and_set_text(shape, "CampusMap结构体封装节点、邻接矩阵、节点数、路径数")
        elif "功能模块个数" in text and len(text) < 10:
            clear_and_set_text(shape, "功能模块个数")
        elif "12个功能模块" in text or "节点/路径增删查" in text:
            clear_and_set_text(shape, "16个函数：节点/路径增删查、最短路径、文件读写、图形化展示")
        elif "邻接矩阵" in text and len(text) < 20:
            clear_and_set_text(shape, "邻接矩阵（二维数组），O(1)随机访问")

# ========== SLIDE 5: Architecture ==========
slide = prs.slides[4]
for shape in slide.shapes:
    if not shape.has_text_frame:
        continue
    text = shape.text_frame.text.strip()
    
    if shape.name == "矩形 4":
        clear_and_set_text(shape, "系统架构图\n\n主菜单控制模块 → 节点管理模块 → 路径管理模块 → 查询模块 → 地图维护模块 → 数据持久化模块\n\n核心交互：主菜单通过while(1)+switch-case调度16个功能函数，各模块共享全局CampusMap结构体")
    elif shape.name == "文本框 6":
        clear_and_set_text(shape, "系统功能架构图\n\n主菜单 → 添加节点/路径 → 查询节点/矩阵/路径 → 删除节点/路径 → 保存/加载地图 → 图形化展示")

# ========== SLIDE 6: Architecture 2 ==========
slide = prs.slides[5]
for shape in slide.shapes:
    if not shape.has_text_frame:
        continue
    if shape.name == "矩形 4":
        clear_and_set_text(shape, "核心数据结构：\n\nstruct CampusMap {\n  char nodes[MAX][50];   // 节点名称\n  int graph[MAX][MAX];   // 邻接矩阵\n  int nodeNum;           // 节点总数\n  int edgeNum;           // 路径总数\n};\n\n#define MAX 100\n#define INF INT_MAX")

# ========== SLIDE 8: Demo ==========
slide = prs.slides[7]
for shape in slide.shapes:
    if not shape.has_text_frame:
        continue
    if shape.name == "矩形 4":
        clear_and_set_text(shape, "系统功能演示与运行效果\n\n1. 启动系统：输入口令1进入主菜单\n2. 添加节点：图书馆、教学楼、食堂、宿舍等校园地标\n3. 添加路径：设置各节点间路径长度\n4. 查询最短路径：Dijkstra算法自动计算最优路线\n5. 查看邻接矩阵：直观展示图结构\n6. 保存/加载地图：数据持久化到campus_map.txt")

# ========== SLIDE 10: Module Overview ==========
slide = prs.slides[9]
for shape in slide.shapes:
    if not shape.has_text_frame:
        continue
    text = shape.text_frame.text.strip()
    
    if shape.name.startswith("Docer"):
        if "模块一" in text and len(text) < 10:
            clear_and_set_text(shape, "模块一")
        elif "模块二" in text and len(text) < 10:
            clear_and_set_text(shape, "模块二")
        elif "模块三" in text and len(text) < 10:
            clear_and_set_text(shape, "模块三")
        elif "模块四" in text and len(text) < 10:
            clear_and_set_text(shape, "模块四")
        elif "插入新记录" in text:
            clear_and_set_text(shape, "添加节点/路径")
        elif "删除记录" in text:
            clear_and_set_text(shape, "删除节点/路径")
        elif "修改记录" in text:
            clear_and_set_text(shape, "查询节点/矩阵")
        elif "显示记录" in text:
            clear_and_set_text(shape, "保存/加载地图")
        elif "节点与路径管理" in text:
            clear_and_set_text(shape, "节点与路径管理")
        elif "地图查询与展示" in text:
            clear_and_set_text(shape, "地图查询与展示")
        elif "最短路径计算" in text:
            clear_and_set_text(shape, "最短路径计算")
        elif "数据持久化" in text:
            clear_and_set_text(shape, "数据持久化")

# ========== SLIDE 11: Algorithm ==========
slide = prs.slides[10]
for shape in slide.shapes:
    if not shape.has_text_frame:
        continue
    text = shape.text_frame.text.strip()
    
    if shape.name == "文本框 1":
        clear_and_set_text(shape, "Dijkstra算法详解")
    elif shape.name == "文本框 5":
        clear_and_set_text(shape, """算法逻辑与复杂度分析（Dijkstra最短路径算法）

算法步骤：
1. 初始化dist、visit、pre三个辅助数组
2. 循环选取未访问节点中距离最小的节点u
3. 对u的所有邻接节点执行松弛操作更新最短路径
4. 重复直至所有节点访问完毕

时间复杂度：O(n²) — 邻接矩阵实现
空间复杂度：O(n) — dist、pre、visit数组

适用条件：非负权图，满足本系统路径长度为正的设定""")

# ========== SLIDE 12: Code Details ==========
slide = prs.slides[11]
for shape in slide.shapes:
    if not shape.has_text_frame:
        continue
    text = shape.text_frame.text.strip()
    
    if shape.name == "文本框 1":
        clear_and_set_text(shape, "关键代码细节解析")
    elif shape.name == "文本框 5":
        clear_and_set_text(shape, """关键代码细节解析（指针/内存/函数调用）

1. 结构体封装：
   struct CampusMap {
     char nodes[MAX][50];
     int graph[MAX][MAX];
     int nodeNum, edgeNum;
   };

2. 邻接矩阵初始化：
   graph[i][j] = (i==j) ? 0 : INF;

3. 路径回溯（递归）：
   void showPath(int end, int pre[]) {
     if (pre[end] != end) {
       showPath(pre[end], pre);
       printf(" -> ");
     }
     printf("%s", nodes[end]);
   }

4. 文件持久化：
   fprintf(fp, "%d\n", nodeNum);
   // 写入节点名称 + 完整邻接矩阵""")

# Save
output_path = 'C:/Users/zero/.openclaw/workspace/校园导航系统_作品演示PPT_修改版.pptx'
prs.save(output_path)
print(f"Saved to: {output_path}")

# Verify
print("\n=== Verification ===")
prs2 = Presentation(output_path)
for i, slide in enumerate(prs2.slides):
    print(f"\n--- Slide {i+1} ---")
    for shape in slide.shapes:
        if shape.has_text_frame:
            text = shape.text_frame.text.strip()
            if text and len(text) > 3:
                print(f"  {shape.name}: {text[:80]}")
